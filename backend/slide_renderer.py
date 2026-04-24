import io
import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image


def _collect_pics(shapes):
    pics = []
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics.append(s)
        elif s.shape_type == MSO_SHAPE_TYPE.GROUP:
            pics.extend(_collect_pics(s.shapes))
    return pics


def render_slides(pptx_path: str, output_dir: str) -> dict:
    """
    Extract screenshots from each PPTX slide and save as PNG files.
    Single-image slides: blob saved directly.
    Multi-image slides: composited onto a bounding-box canvas.
    Returns {pptx_slide_number (1-based): filename}.
    """
    prs = Presentation(pptx_path)
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    os.makedirs(output_dir, exist_ok=True)
    rendered = {}

    for i, slide in enumerate(prs.slides):
        pptx_num = i + 1
        pics = _collect_pics(slide.shapes)
        if not pics:
            continue

        filename = f"slide_{pptx_num:02d}.png"
        out_path = os.path.join(output_dir, filename)

        if os.path.exists(out_path):
            rendered[pptx_num] = filename
            continue

        if len(pics) == 1:
            img = Image.open(io.BytesIO(pics[0].image.blob)).convert("RGB")
            img.save(out_path)
        else:
            left_min = min(p.left for p in pics)
            top_min = min(p.top for p in pics)
            right_max = max(p.left + p.width for p in pics)
            bottom_max = max(p.top + p.height for p in pics)
            bbox_w = right_max - left_min
            bbox_h = bottom_max - top_min

            canvas_w = 1280
            canvas_h = int(canvas_w * bbox_h / bbox_w)
            canvas = Image.new("RGBA", (canvas_w, canvas_h), (255, 255, 255, 255))

            # Draw largest (base) image first, then overlays on top
            for p in sorted(pics, key=lambda x: -len(x.image.blob)):
                img = Image.open(io.BytesIO(p.image.blob)).convert("RGBA")
                px = int((p.left - left_min) / bbox_w * canvas_w)
                py = int((p.top - top_min) / bbox_h * canvas_h)
                pw = int(p.width / bbox_w * canvas_w)
                ph = int(p.height / bbox_h * canvas_h)
                img = img.resize((pw, ph), Image.LANCZOS)
                canvas.paste(img, (px, py), img)

            canvas.convert("RGB").save(out_path)

        rendered[pptx_num] = filename

    return rendered


def build_slide_map(pptx_path: str) -> dict:
    """
    Scan PPTX slides for text labels 'Slide N' and 'Sim X'.
    Returns {(sim_key: str, slide_num: int) -> pptx_slide_number: int}.
    """
    prs = Presentation(pptx_path)
    result = {}

    for i, slide in enumerate(prs.slides):
        pptx_num = i + 1
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]

        slide_num = None
        sim_key = None

        for text in texts:
            m = re.match(r"Slide\s+(\d+)$", text, re.IGNORECASE)
            if m:
                slide_num = int(m.group(1))
            lower = text.lower()
            if "sim 1" in lower or "sim1" in lower:
                sim_key = "sim-1"
            elif "sim 2" in lower or "sim2" in lower:
                sim_key = "sim-2"

        if sim_key and slide_num:
            result[(sim_key, slide_num)] = pptx_num

    return result
